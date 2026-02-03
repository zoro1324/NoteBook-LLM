"""
Document Processing Service
Handles extraction and processing of various document types.
"""
import os
import re
from pathlib import Path
from typing import List, Tuple
from django.core.files.uploadedfile import UploadedFile


class DocumentProcessor:
    """Base class for document processing"""
    
    @staticmethod
    def detect_file_type(filename: str) -> str:
        """Detect file type from extension"""
        ext = Path(filename).suffix.lower()
        
        mapping = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.doc': 'docx',
            '.txt': 'txt',
            '.md': 'md',
            '.mp3': 'audio',
            '.wav': 'audio',
            '.m4a': 'audio',
            '.mp4': 'video',
            '.avi': 'video',
            '.mov': 'video',
        }
        
        return mapping.get(ext, 'txt')
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 512, overlap: int = 50) -> List[Tuple[str, int, int]]:
        """
        Split text into overlapping chunks for embeddings.
        
        Returns:
            List of (chunk_text, start_char, end_char) tuples
        """
        chunks = []
        words = text.split()
        
        if not words:
            return chunks
        
        current_chunk = []
        current_size = 0
        char_position = 0
        
        for word in words:
            word_size = len(word) + 1  # +1 for space
            
            if current_size + word_size > chunk_size and current_chunk:
                # Save current chunk
                chunk_text = ' '.join(current_chunk)
                start_pos = char_position - len(chunk_text)
                chunks.append((chunk_text, start_pos, char_position))
                
                # Start new chunk with overlap
                overlap_words = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_words + [word]
                current_size = sum(len(w) + 1 for w in current_chunk)
                char_position += word_size
            else:
                current_chunk.append(word)
                current_size += word_size
                char_position += word_size
        
        # Add final chunk
        if current_chunk:
            chunk_text = ' '.join(current_chunk)
            start_pos = char_position - len(chunk_text)
            chunks.append((chunk_text, start_pos, char_position))
        
        return chunks


class PDFProcessor(DocumentProcessor):
    """
    Process PDF documents using OCR as PRIMARY method.
    Uses doctr (deep learning OCR) for best accuracy on all PDF types.
    """
    
    # Cache the OCR model to avoid reloading
    _ocr_model = None
    
    @classmethod
    def get_ocr_model(cls):
        """Get or create the OCR model (cached for performance)"""
        if cls._ocr_model is None:
            try:
                from doctr.models import ocr_predictor
                print("[PDF] Loading doctr OCR model...")
                cls._ocr_model = ocr_predictor(pretrained=True)
                print("[PDF] OCR model loaded successfully!")
            except ImportError:
                print("[PDF] doctr not installed, OCR unavailable")
                return None
            except Exception as e:
                print(f"[PDF] Failed to load OCR model: {e}")
                return None
        return cls._ocr_model
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, dict]:
        """
        Extract text from PDF using OCR as primary method.
        
        Returns:
            (extracted_text, metadata)
        """
        metadata = {'pages': 0, 'method': 'unknown'}
        
        # Strategy 1: Use doctr OCR (PRIMARY - works for all PDFs)
        text, meta = PDFProcessor._extract_with_doctr_ocr(file_path)
        if text and len(text.strip()) > 20:
            metadata.update(meta)
            metadata['method'] = 'doctr_ocr'
            return text, metadata
        
        # Strategy 2: Fallback to PyMuPDF
        text, meta = PDFProcessor._extract_with_pymupdf(file_path)
        if text and len(text.strip()) > 20:
            metadata.update(meta)
            metadata['method'] = 'pymupdf'
            return text, metadata
        
        # Strategy 3: Fallback to PyPDF2
        text, meta = PDFProcessor._extract_with_pypdf2(file_path)
        if text and len(text.strip()) > 20:
            metadata.update(meta)
            metadata['method'] = 'pypdf2'
            return text, metadata
        
        # If all methods failed
        if not text:
            raise Exception("PDF extraction failed: No text could be extracted from the PDF.")
        
        return text, metadata
    
    @staticmethod
    def _extract_with_doctr_ocr(file_path: str) -> Tuple[str, dict]:
        """Extract text using doctr deep learning OCR - PRIMARY method"""
        try:
            from doctr.io import DocumentFile
            
            print(f"[PDF] Extracting text with doctr OCR: {file_path}")
            
            # Get or load the model
            model = PDFProcessor.get_ocr_model()
            if model is None:
                return "", {'pages': 0, 'error': 'OCR model not available'}
            
            # Load PDF as images and run OCR
            doc = DocumentFile.from_pdf(file_path)
            result = model(doc)
            
            # Extract text from OCR result
            text_pages = []
            for page_idx, page in enumerate(result.pages):
                page_text = []
                for block in page.blocks:
                    block_lines = []
                    for line in block.lines:
                        line_text = " ".join([word.value for word in line.words])
                        block_lines.append(line_text)
                    page_text.append(" ".join(block_lines))
                
                if page_text:
                    page_content = "\n".join(page_text)
                    text_pages.append(f"--- Page {page_idx + 1} ---\n{page_content}")
            
            full_text = '\n\n'.join(text_pages)
            word_count = len(full_text.split())
            print(f"[PDF] doctr OCR extracted {word_count} words from {len(result.pages)} pages")
            
            return full_text, {'pages': len(result.pages), 'ocr_engine': 'doctr', 'word_count': word_count}
            
        except ImportError as e:
            print(f"[PDF] doctr not installed: {e}")
            return "", {'pages': 0, 'error': 'doctr not installed'}
        except Exception as e:
            print(f"[PDF] doctr OCR failed: {e}")
            return "", {'pages': 0, 'error': str(e)}
    
    @staticmethod
    def _extract_with_pymupdf(file_path: str) -> Tuple[str, dict]:
        """Extract text using PyMuPDF (fitz) - fallback method"""
        try:
            import fitz  # PyMuPDF
            
            print(f"[PDF] Trying PyMuPDF extraction: {file_path}")
            
            text_pages = []
            metadata = {'pages': 0}
            
            doc = fitz.open(file_path)
            metadata['pages'] = len(doc)
            
            for page_num, page in enumerate(doc):
                text = page.get_text("text")
                if text.strip():
                    text_pages.append(f"--- Page {page_num + 1} ---\n{text}")
            
            doc.close()
            full_text = '\n\n'.join(text_pages)
            print(f"[PDF] PyMuPDF extracted {len(full_text.split())} words")
            return full_text, metadata
            
        except ImportError:
            return "", {'pages': 0, 'error': 'PyMuPDF not installed'}
        except Exception as e:
            print(f"[PDF] PyMuPDF failed: {e}")
            return "", {'pages': 0, 'error': str(e)}
    
    @staticmethod
    def _extract_with_pypdf2(file_path: str) -> Tuple[str, dict]:
        """Extract text using PyPDF2 - last resort fallback"""
        try:
            import PyPDF2
            
            print(f"[PDF] Trying PyPDF2 extraction: {file_path}")
            
            text_pages = []
            metadata = {'pages': 0}
            
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                metadata['pages'] = len(reader.pages)
                
                for page_num, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        text_pages.append(f"--- Page {page_num + 1} ---\n{text}")
            
            full_text = '\n\n'.join(text_pages)
            print(f"[PDF] PyPDF2 extracted {len(full_text.split())} words")
            return full_text, metadata
            
        except Exception as e:
            print(f"[PDF] PyPDF2 failed: {e}")
            return "", {'pages': 0, 'error': str(e)}


class TextProcessor(DocumentProcessor):
    """Process plain text and markdown files"""
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, dict]:
        """Extract text from text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            metadata = {
                'lines': len(text.split('\n')),
                'chars': len(text),
            }
            
            return text, metadata
        except Exception as e:
            raise Exception(f"Text extraction failed: {str(e)}")


class DocxProcessor(DocumentProcessor):
    """Process Word documents with table support"""
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, dict]:
        """Extract text and tables from DOCX"""
        try:
            import docx
            
            doc = docx.Document(file_path)
            
            text_parts = []
            table_count = 0
            
            for element in doc.element.body:
                # Handle paragraphs
                if element.tag.endswith('p'):
                    for para in doc.paragraphs:
                        if para._element == element and para.text.strip():
                            # Check for headings
                            if para.style and para.style.name.startswith('Heading'):
                                level = para.style.name[-1] if para.style.name[-1].isdigit() else '2'
                                text_parts.append(f"{'#' * int(level)} {para.text}")
                            else:
                                text_parts.append(para.text)
                            break
                
                # Handle tables
                elif element.tag.endswith('tbl'):
                    for table in doc.tables:
                        if table._tbl == element:
                            table_count += 1
                            table_md = DocxProcessor._table_to_markdown(table)
                            text_parts.append(f"\n{table_md}\n")
                            break
            
            full_text = '\n\n'.join(text_parts)
            
            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': table_count,
                'chars': len(full_text),
            }
            
            return full_text, metadata
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    @staticmethod
    def _table_to_markdown(table) -> str:
        """Convert DOCX table to markdown format"""
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
            
            # Add header separator after first row
            if i == 0:
                rows.append('|' + '|'.join(['---'] * len(cells)) + '|')
        
        return '\n'.join(rows)


class PPTXProcessor(DocumentProcessor):
    """Process PowerPoint presentations"""
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, dict]:
        """Extract text from PPTX slides"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_text = []
            total_shapes = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"## Slide {slide_num}")
                
                for shape in slide.shapes:
                    total_shapes += 1
                    
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # Check if it's a title
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            if shape.placeholder_format.type == 1:  # Title
                                slide_content.append(f"### {text}")
                            else:
                                slide_content.append(text)
                        else:
                            slide_content.append(text)
                    
                    # Handle tables in slides
                    if shape.has_table:
                        table_md = PPTXProcessor._table_to_markdown(shape.table)
                        slide_content.append(f"\n{table_md}\n")
                
                if len(slide_content) > 1:  # More than just the slide header
                    slides_text.append('\n\n'.join(slide_content))
            
            full_text = '\n\n---\n\n'.join(slides_text)
            
            metadata = {
                'slides': len(prs.slides),
                'shapes': total_shapes,
                'chars': len(full_text),
            }
            
            return full_text, metadata
            
        except ImportError:
            raise Exception("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    @staticmethod
    def _table_to_markdown(table) -> str:
        """Convert PPTX table to markdown format"""
        rows = []
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('|', '\\|')
                cells.append(cell_text)
            rows.append('| ' + ' | '.join(cells) + ' |')
            
            if i == 0:
                rows.append('|' + '|'.join(['---'] * len(cells)) + '|')
        
        return '\n'.join(rows)


class ImageProcessor(DocumentProcessor):
    """Process image files with OCR"""
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, dict]:
        """Extract text from image using OCR"""
        try:
            # Try using doctr first (same as PDF processor)
            from doctr.io import DocumentFile
            from doctr.models import ocr_predictor
            
            print(f"[Image] Processing with doctr OCR: {file_path}")
            
            # Load image
            doc = DocumentFile.from_images(file_path)
            
            # Get OCR model
            model = ocr_predictor(pretrained=True)
            result = model(doc)
            
            # Extract text
            text_lines = []
            for page in result.pages:
                for block in page.blocks:
                    for line in block.lines:
                        line_text = " ".join([word.value for word in line.words])
                        text_lines.append(line_text)
            
            full_text = '\n'.join(text_lines)
            
            metadata = {
                'source': 'doctr_ocr',
                'chars': len(full_text),
            }
            
            return full_text, metadata
            
        except ImportError:
            # Fallback to basic PIL text detection message
            print("[Image] doctr not available, returning placeholder")
            return "Image content requires OCR processing. Install doctr for OCR support.", {'source': 'placeholder'}
        except Exception as e:
            raise Exception(f"Image OCR failed: {str(e)}")


class AudioProcessor(DocumentProcessor):
    """Process audio files with transcription"""
    
    @staticmethod
    def extract_text(file_path: str) -> Tuple[str, dict]:
        """Transcribe audio file using Faster Whisper"""
        try:
            from faster_whisper import WhisperModel
            from django.conf import settings
            
            config = settings.APP_CONFIG.get('transcription', {})
            model_size = config.get('model', 'tiny')
            device = config.get('device', 'cpu')  # Default to CPU for compatibility
            
            model = WhisperModel(model_size, device=device)
            segments, info = model.transcribe(file_path)
            
            text_segments = []
            for segment in segments:
                text_segments.append(segment.text)
            
            full_text = ' '.join(text_segments)
            
            metadata = {
                'language': info.language,
                'duration': info.duration,
            }
            
            return full_text, metadata
        except Exception as e:
            raise Exception(f"Audio transcription failed: {str(e)}")


# Factory function to get appropriate processor
def get_processor(file_type: str):
    """Get processor instance for file type"""
    processors = {
        'pdf': PDFProcessor,
        'txt': TextProcessor,
        'md': TextProcessor,
        'docx': DocxProcessor,
        'pptx': PPTXProcessor,
        'image': ImageProcessor,
        'audio': AudioProcessor,
    }
    
    return processors.get(file_type, TextProcessor)()

